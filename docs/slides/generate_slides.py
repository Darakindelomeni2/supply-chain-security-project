#!/usr/bin/env python3
"""
Génère le PowerPoint d'introduction (Jour 1) du projet
"Sécurité de la chaîne d'approvisionnement logicielle (SLSA)".

Usage :
    pip install python-pptx
    python generate_slides.py
    # -> produit : intro-supply-chain-security.pptx  (ignoré par git)

Le .pptx est un artefact binaire régénérable : on versionne CE script, pas le .pptx.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ----------------------------------------------------------------------------- palette
NAVY   = RGBColor(0x0B, 0x1F, 0x3A)   # fond sombre / titres
NAVY2  = RGBColor(0x13, 0x2A, 0x4D)
TEAL   = RGBColor(0x10, 0xB9, 0x81)   # accent "confiance / vérifié" (vert émeraude)
TEALD  = RGBColor(0x0B, 0x81, 0x5A)
RED    = RGBColor(0xEF, 0x44, 0x44)   # accent "attaque / bloqué"
AMBER  = RGBColor(0xF5, 0x9E, 0x0B)
INK    = RGBColor(0x1F, 0x29, 0x37)   # texte foncé
GREY   = RGBColor(0x6B, 0x72, 0x80)
LIGHT  = RGBColor(0xF3, 0xF4, 0xF6)   # fond clair
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
CARD   = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"
W, H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


# ----------------------------------------------------------------------------- helpers
def slide():
    return prs.slides.add_slide(BLANK)


def bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


def rect(s, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE, line=None):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1.25)
    sp.shadow.inherit = False
    return sp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6):
    """runs = liste de (texte, taille, couleur, gras) OU liste de paragraphes (liste de runs)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if runs and isinstance(runs[0], tuple):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        for (t, sz, col, bold) in para:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(sz)
            r.font.color.rgb = col
            r.font.bold = bold
            r.font.name = FONT
    return tb


def bullets(s, x, y, w, h, items, size=18, color=INK, gap=10):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        lvl = 0
        mark, txt = "▸  ", it
        if isinstance(it, tuple):     # (niveau, texte)
            lvl, txt = it
            mark = "▸  " if lvl == 0 else "–  "
        r = p.add_run(); r.text = mark + txt
        r.font.size = Pt(size if lvl == 0 else size - 3)
        r.font.color.rgb = color if lvl == 0 else GREY
        r.font.name = FONT
        r.font.bold = False
    return tb


def header(s, kicker, title, accent=TEAL):
    """En-tête standard des slides de contenu (fond clair)."""
    bg(s, WHITE)
    rect(s, 0, 0, Inches(0.28), H, accent)                       # barre latérale accent
    text(s, Inches(0.6), Inches(0.45), Inches(12), Inches(0.4),
         [(kicker.upper(), 13, accent, True)])
    text(s, Inches(0.6), Inches(0.78), Inches(12.2), Inches(0.9),
         [(title, 30, NAVY, True)])
    rect(s, Inches(0.62), Inches(1.62), Inches(1.7), Pt(3), accent)


def footer(s, n, dark=False):
    col = RGBColor(0xB8, 0xC2, 0xD0) if dark else GREY
    text(s, Inches(0.6), Inches(7.02), Inches(9), Inches(0.35),
         [("Projet 5ᵉ année — Sécurité de la chaîne d'approvisionnement logicielle (SLSA)", 9, col, False)])
    text(s, Inches(11.6), Inches(7.02), Inches(1.4), Inches(0.35),
         [(str(n), 9, col, False)], align=PP_ALIGN.RIGHT)


def chip(s, x, y, w, label, color):
    c = rect(s, x, y, w, Inches(0.42), color, MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = c.text_frame; tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
    return c


def arrow(s, x, y):
    a = rect(s, x, y, Inches(0.32), Inches(0.42), GREY, MSO_SHAPE.CHEVRON)
    return a


# ============================================================================= 1. TITRE
s = slide(); bg(s, NAVY)
rect(s, 0, 0, W, Inches(0.16), TEAL)
rect(s, 0, Inches(7.34), W, Inches(0.16), TEAL)
text(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.5),
     [("PROJET TECHNIQUE · 5ᵉ ANNÉE DEVOPS / CLOUD / INFRA", 15, TEAL, True)])
text(s, Inches(0.9), Inches(2.15), Inches(11.6), Inches(2.2),
     [[("Sécuriser la chaîne", 48, WHITE, True)],
      [("d'approvisionnement logicielle ", 48, WHITE, True), ("(SLSA)", 48, TEAL, True)]])
text(s, Inches(0.9), Inches(4.35), Inches(11.4), Inches(1.0),
     [("« Comment prouver que l'image qui tourne en production est bien celle que "
       "NOUS avons construite — et pas une version piégée ? »", 19, RGBColor(0xC7,0xD2,0xE0), False)])
chip(s, Inches(0.9), Inches(5.7), Inches(3.3), "3 jours · 1,5 j de labs", TEALD)
chip(s, Inches(4.4), Inches(5.7), Inches(2.2), "QCM individuel", NAVY2)
chip(s, Inches(6.8), Inches(5.7), Inches(2.6), "Soutenance + démo", NAVY2)
text(s, Inches(0.9), Inches(6.5), Inches(11), Inches(0.4),
     [("Jour 1 — Introduction", 14, RGBColor(0x9F,0xB0,0xC3), True)])

# ============================================================================= 2. PROBLÈME
s = slide(); header(s, "Le point de départ", "Votre pipeline CI/CD est-il digne de confiance ?")
text(s, Inches(0.6), Inches(1.95), Inches(12), Inches(0.6),
     [("Vous savez tous faire : build → test → scan → deploy. Mais une fois l'image en production…", 18, INK, False)])
q = rect(s, Inches(0.6), Inches(2.7), Inches(12.1), Inches(2.4), LIGHT, MSO_SHAPE.ROUNDED_RECTANGLE)
bullets(s, Inches(1.0), Inches(2.95), Inches(11.4), Inches(2.0), [
    "Le registry peut être compromis → l'image est remplacée.",
    "Un accès au cluster permet de déployer N'IMPORTE quelle image.",
    "Une dépendance peut cacher une backdoor (cf. XZ Utils, 2024).",
    "Un runner de CI compromis peut injecter du code DANS le build (cf. SolarWinds).",
], size=17, gap=9)
text(s, Inches(0.6), Inches(5.4), Inches(12), Inches(1.0),
     [[("Un ", 19, INK, False), ("docker pull", 19, RED, True),
       (" ne vérifie RIEN de tout cela. Et « le scan Trivy était vert » ne prouve pas que "
        "l'image DÉPLOYÉE est celle qui a été scannée.", 19, INK, False)]])
footer(s, 2)

# ============================================================================= 3. ATTAQUES RÉELLES
s = slide(); header(s, "Pourquoi ça compte", "Les attaques ne visent plus votre app — mais votre chaîne", RED)
cards = [
    ("SolarWinds — 2020", "Code malveillant injecté dans le BUILD, signé par l'éditeur, poussé à 18 000 clients."),
    ("Codecov — 2021", "Script CI modifié : exfiltration des secrets des pipelines de milliers de projets."),
    ("Dependency confusion — 2021", "Faux paquets « internes » publiés sur les registries publics, installés par erreur."),
    ("XZ Utils / liblzma — 2024", "Backdoor introduite sur 3 ANS dans une dépendance open source très répandue."),
]
x0, y0, cw, ch, gx, gy = Inches(0.6), Inches(2.05), Inches(5.95), Inches(1.95), Inches(0.2), Inches(0.25)
for i, (t, d) in enumerate(cards):
    cx = x0 + (cw + gx) * (i % 2)
    cy = y0 + (ch + gy) * (i // 2)
    rect(s, cx, cy, cw, ch, LIGHT, MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, cx, cy, Inches(0.12), ch, RED)
    text(s, cx + Inches(0.35), cy + Inches(0.18), cw - Inches(0.55), Inches(0.5),
         [(t, 17, NAVY, True)])
    text(s, cx + Inches(0.35), cy + Inches(0.72), cw - Inches(0.55), Inches(1.1),
         [(d, 14, INK, False)])
footer(s, 3)

# ============================================================================= 4. LA CHAÎNE VÉRIFIABLE
s = slide(); header(s, "La réponse", "Transformer le pipeline en chaîne VÉRIFIABLE")
steps = [("Code", TEAL), ("Build", TEAL), ("SBOM", TEAL), ("Scan", AMBER),
         ("Sign", TEALD), ("Attest", TEALD), ("Registry", NAVY2)]
x = Inches(0.6); y = Inches(2.35); cw = Inches(1.42)
for i, (lbl, col) in enumerate(steps):
    chip(s, x, y, cw, lbl, col)
    x = x + cw + Inches(0.05)
    if i < len(steps) - 1:
        arrow(s, x, y); x = x + Inches(0.37)
text(s, Inches(0.6), Inches(2.92), Inches(12), Inches(0.4),
     [("SBOM (Syft) → scan (Grype, casse si CVE critique) → signature (cosign) → attestations SBOM + provenance", 13, GREY, False)])
# le cluster
rect(s, Inches(0.6), Inches(3.55), Inches(12.1), Inches(2.2), NAVY, MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, Inches(0.95), Inches(3.75), Inches(11.4), Inches(0.5),
     [[("Cluster Kubernetes (kind/k3s) + ", 18, WHITE, True), ("KYVERNO", 18, TEAL, True),
       ("  — admission control", 18, WHITE, True)]])
checks = [
    "Signée par NOTRE identité ?", "Attestation de provenance ?",
    "Registry autorisé + par digest ?", "Pas de CVE critique / pas de :latest ?",
]
cx = Inches(0.95)
for i, c in enumerate(checks):
    bx = Inches(0.95) + (Inches(2.95)) * i
    rect(s, bx, Inches(4.45), Inches(2.75), Inches(1.05), NAVY2, MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, bx + Inches(0.15), Inches(4.55), Inches(2.5), Inches(0.9),
         [[("✓ ", 14, TEAL, True), (c, 13, RGBColor(0xD7,0xE0,0xEC), False)]])
text(s, Inches(0.6), Inches(5.95), Inches(12), Inches(0.8),
     [[("Sinon → ", 19, INK, False), ("❌ REFUSÉE", 22, RED, True),
       ("   ·   Tout est vérifié → ", 19, INK, False), ("✅ le pod tourne", 20, TEALD, True)]])
footer(s, 4)

# ============================================================================= 5. LES 4 BRIQUES
s = slide(); header(s, "Ce que vous allez mettre en œuvre", "Les 4 briques de la chaîne")
bricks = [
    ("1 · SBOM", "Syft + Grype", "L'« étiquette de composition » de l'image : tous les paquets et versions. « Suis-je affecté par la CVE du jour ? » en secondes.", TEAL),
    ("2 · Signature", "cosign / Sigstore", "Preuve cryptographique « c'est bien NOUS ». Mode keyless (identité OIDC) journalisé dans Rekor : aucune clé à gérer.", TEALD),
    ("3 · Attestations", "provenance SLSA", "Affirmations signées attachées à l'image : le SBOM lui-même + la provenance (qui/quoi/d'où/quand du build).", AMBER),
    ("4 · Admission", "Kyverno", "Le gardien du cluster : vérifie signature + attestations à CHAQUE déploiement et REFUSE l'inconnu.", RED),
]
x0, y0, cw, ch, gx = Inches(0.6), Inches(2.05), Inches(2.93), Inches(4.4), Inches(0.15)
for i, (t, sub, d, col) in enumerate(bricks):
    cx = x0 + (cw + gx) * i
    rect(s, cx, y0, cw, ch, LIGHT, MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, cx, y0, cw, Inches(0.9), col, MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, cx + Inches(0.2), y0 + Inches(0.12), cw - Inches(0.4), Inches(0.8),
         [[(t, 17, WHITE, True)], [(sub, 12, WHITE, False)]])
    text(s, cx + Inches(0.22), y0 + Inches(1.15), cw - Inches(0.44), Inches(3.0),
         [(d, 14, INK, False)])
footer(s, 5)

# ============================================================================= 6. SLSA
s = slide(); header(s, "Le référentiel", "SLSA — des niveaux de garantie sur la provenance", TEALD)
text(s, Inches(0.6), Inches(1.95), Inches(12), Inches(0.5),
     [[("SLSA", 18, TEALD, True),
       (" (Supply-chain Levels for Software Artifacts, OpenSSF) — « salsa ». On visera ", 17, INK, False),
       ("L2", 18, TEALD, True), (".", 17, INK, False)]])
levels = [
    ("L1", "La provenance EXISTE", "Le build enregistre comment l'artefact a été fait.", RGBColor(0xCF,0xE9,0xDD)),
    ("L2", "Provenance SIGNÉE + build hébergé", "Build sur une plateforme (pas sur un poste), provenance signée. ← notre cible", TEAL),
    ("L3", "Build renforcé & isolé", "Isolation forte, provenance infalsifiable, non contournable.", RGBColor(0xCF,0xE9,0xDD)),
]
y = Inches(2.75)
for lvl, t, d, col in levels:
    rect(s, Inches(0.6), y, Inches(12.1), Inches(1.15), LIGHT, MSO_SHAPE.ROUNDED_RECTANGLE)
    badge = rect(s, Inches(0.8), y + Inches(0.2), Inches(1.15), Inches(0.75), col, MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = badge.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = lvl; r.font.size = Pt(24); r.font.bold = True
    r.font.color.rgb = NAVY; r.font.name = FONT
    text(s, Inches(2.2), y + Inches(0.16), Inches(10.2), Inches(0.5), [(t, 18, NAVY, True)])
    text(s, Inches(2.2), y + Inches(0.62), Inches(10.2), Inches(0.5), [(d, 14, INK, False)])
    y = y + Inches(1.3)
footer(s, 6)

# ============================================================================= 7. LA DÉMO CIBLE
s = slide(); header(s, "L'objectif concret", "La démo de soutenance : attaque / défense", RED)
# colonne OK
rect(s, Inches(0.6), Inches(2.05), Inches(5.9), Inches(4.3), RGBColor(0xE7,0xF7,0xF0), MSO_SHAPE.ROUNDED_RECTANGLE)
rect(s, Inches(0.6), Inches(2.05), Inches(5.9), Inches(0.85), TEALD, MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, Inches(0.9), Inches(2.2), Inches(5.3), Inches(0.6), [[("✅  Image légitime", 20, WHITE, True)]])
bullets(s, Inches(0.95), Inches(3.1), Inches(5.3), Inches(3.1), [
    "Signée par notre identité", "SBOM + provenance attachés",
    "Registry autorisé, par digest", "→ le pod démarre, l'app répond",
], size=16, color=INK, gap=12)
# colonne KO
rect(s, Inches(6.8), Inches(2.05), Inches(5.9), Inches(4.3), RGBColor(0xFB,0xE9,0xE9), MSO_SHAPE.ROUNDED_RECTANGLE)
rect(s, Inches(6.8), Inches(2.05), Inches(5.9), Inches(0.85), RED, MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, Inches(7.1), Inches(2.2), Inches(5.3), Inches(0.6), [[("❌  Image piégée / non signée", 20, WHITE, True)]])
bullets(s, Inches(7.15), Inches(3.1), Inches(5.3), Inches(3.1), [
    "Non signée → refusée", "Modifiée après signature → refusée",
    "Mauvais registry ou :latest → refusée", "→ Kyverno BLOQUE, en direct",
], size=16, color=INK, gap=12)
footer(s, 7)

# ============================================================================= 8. PLANNING
s = slide(); header(s, "Organisation", "3 jours : ~1,5 j de projet, puis QCM & soutenances")
rows = [
    ("Jour 1", "Chaîne vérifiable", "Labs 0→2 : build, SBOM, scan, signature, attestations", TEAL),
    ("Jour 2 matin", "Le cluster qui refuse", "Labs 3→4 : Kyverno + attaque/défense  ·  FIN du projet (1,5 j)", TEALD),
    ("Jour 2 aprèm", "Intégration & rédaction", "Lab 5 (CI GitHub Actions, bonus) + rapport + threat model", NAVY2),
    ("Jour 3 matin", "QCM individuel", "Répétition démo + QCM (25-30 min)", AMBER),
    ("Jour 3 aprèm", "Soutenances", "12 min démo/présentation + 5 min Q/R par groupe", RED),
]
y = Inches(2.0)
for d, t, desc, col in rows:
    rect(s, Inches(0.6), y, Inches(12.1), Inches(0.92), LIGHT, MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, Inches(0.6), y, Inches(2.4), Inches(0.92), col, MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, Inches(0.7), y + Inches(0.22), Inches(2.3), Inches(0.5), [(d, 15, WHITE, True)], align=PP_ALIGN.CENTER)
    text(s, Inches(3.2), y + Inches(0.1), Inches(9.3), Inches(0.45), [(t, 16, NAVY, True)])
    text(s, Inches(3.2), y + Inches(0.5), Inches(9.3), Inches(0.4), [(desc, 13, INK, False)])
    y = y + Inches(1.02)
footer(s, 8)

# ============================================================================= 9. ÉVALUATION
s = slide(); header(s, "Évaluation", "Comment vous êtes noté·e", TEALD)
data = [("POC & démo", 35, TEAL), ("Rapport + threat model", 25, TEALD),
        ("Soutenance", 20, NAVY2), ("QCM individuel", 20, AMBER)]
y = Inches(2.3)
maxw = Inches(8.2)
for lbl, pct, col in data:
    text(s, Inches(0.6), y, Inches(3.6), Inches(0.5), [(lbl, 17, NInk := INK, True)])
    barw = Emu(int(maxw * (pct / 35)))
    rect(s, Inches(4.3), y + Inches(0.03), Emu(int(maxw)), Inches(0.42), LIGHT, MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, Inches(4.3), y + Inches(0.03), barw, Inches(0.42), col, MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, Inches(12.5), y, Inches(0.7), Inches(0.5), [(f"{pct}%", 16, NAVY, True)], align=PP_ALIGN.RIGHT)
    y = y + Inches(0.75)
text(s, Inches(0.6), Inches(5.7), Inches(12), Inches(1.2),
     [[("Le POC/démo et le QCM pèsent le plus. ", 16, INK, True),
       ("La démo doit montrer un blocage RÉEL et reproductible ; le QCM est individuel.", 16, INK, False)]])
footer(s, 9)

# ============================================================================= 10. LES LABS
s = slide(); header(s, "Le parcours", "Les 5 labs (le cœur des 1,5 jour)")
labs = [
    ("Lab 0", "Setup & première image", "Outils, fork, build, push sur GHCR, notion de digest."),
    ("Lab 1", "SBOM & scan qui casse", "Syft (SBOM) + Grype (gate qui stoppe sur CVE critique)."),
    ("Lab 2", "Signer & attester", "cosign : signature + attestations SBOM et provenance."),
    ("Lab 3", "Cluster qui refuse", "kind + Kyverno : exiger signature, provenance, registry."),
    ("Lab 4", "Attaque / défense", "Image non signée / modifiée → BLOQUÉE. Captures pour la démo."),
    ("Lab 5", "CI de bout en bout (bonus)", "Tout automatiser en GitHub Actions, keyless → vers SLSA L2."),
]
x0, y0, cw, ch, gx, gy = Inches(0.6), Inches(2.05), Inches(3.95), Inches(1.55), Inches(0.13), Inches(0.18)
for i, (t, tt, d) in enumerate(labs):
    cx = x0 + (cw + gx) * (i % 3)
    cy = y0 + (ch + gy) * (i // 3)
    rect(s, cx, cy, cw, ch, LIGHT, MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, cx, cy, Inches(0.12), ch, TEAL)
    text(s, cx + Inches(0.3), cy + Inches(0.12), cw - Inches(0.5), Inches(0.4),
         [[(t + " · ", 14, TEALD, True), (tt, 14, NAVY, True)]])
    text(s, cx + Inches(0.3), cy + Inches(0.6), cw - Inches(0.5), Inches(0.9),
         [(d, 12.5, INK, False)])
footer(s, 10)

# ============================================================================= 11. PRÉREQUIS / CODE FOURNI
s = slide(); header(s, "Avant de commencer", "Prérequis & code fourni", NAVY2)
# gauche : outils
rect(s, Inches(0.6), Inches(2.05), Inches(5.9), Inches(4.5), LIGHT, MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, Inches(0.9), Inches(2.2), Inches(5.3), Inches(0.5), [("🧰  Outils à installer", 18, NAVY, True)])
bullets(s, Inches(0.95), Inches(2.8), Inches(5.3), Inches(3.6), [
    "Docker · kind (ou k3s) · kubectl",
    "Syft · Grype · cosign",
    "git · jq · compte GitHub (+ PAT write:packages)",
    "Kyverno s'installe dans le cluster (Lab 3)",
    "100 % local — aucun cloud requis (Azure en option)",
], size=15, gap=11)
# droite : fourni vs à produire
rect(s, Inches(6.8), Inches(2.05), Inches(5.9), Inches(4.5), RGBColor(0xE7,0xF7,0xF0), MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, Inches(7.1), Inches(2.2), Inches(5.3), Inches(0.5), [("📦  Vous ne partez PAS de zéro", 18, TEALD, True)])
bullets(s, Inches(7.15), Inches(2.8), Inches(5.3), Inches(3.6), [
    "Fourni : app Flask, Dockerfile, 4 politiques Kyverno,",
    (1, "manifeste K8s, workflow CI, tous les labs"),
    "À produire : SBOM, scan, signature, attestations",
    "À personnaliser : remplacer <votre-user> sur votre fork",
    "Repo à forker : github.com/aubinaso/…-security-project",
], size=15, gap=10)
footer(s, 11)

# ============================================================================= 12. À VOUS
s = slide(); bg(s, NAVY)
rect(s, 0, Inches(3.3), W, Inches(0.06), TEAL)
text(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(1.2),
     [[("À vous de jouer 🔐", 40, WHITE, True)]])
text(s, Inches(0.9), Inches(3.6), Inches(11.5), Inches(1.0),
     [[("On ne fait pas confiance — ", 22, RGBColor(0xC7,0xD2,0xE0), False),
       ("on vérifie.", 22, TEAL, True)]])
bullets(s, Inches(0.95), Inches(4.5), Inches(11.4), Inches(2.0), [
    "Constituez vos groupes (2 à 4) et forkez le dépôt.",
    "Ouvrez docs/00-presentation-projet.md puis labs/lab0-setup.md.",
    "Objectif fin J1 : une image SIGNÉE, vérifiable, dans votre registry.",
], size=18, color=RGBColor(0xD7,0xE0,0xEC), gap=12)
footer(s, 12, dark=True)

# ----------------------------------------------------------------------------- save
import os
out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "intro-supply-chain-security.pptx")
prs.save(out)
print("OK ->", out, f"({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
